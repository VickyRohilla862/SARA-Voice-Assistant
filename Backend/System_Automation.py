import os
import subprocess
import webbrowser
import platform
import shutil
import json
from pathlib import Path
from datetime import datetime
import socket

import cv2
import numpy as np
import pyautogui
import pyaudio
import wave
import threading
import keyboard
from datetime import datetime
import pyaudiowpatch

import psutil
import requests
import pyautogui

# ===== AUDIO =====
from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
from comtypes import CLSCTX_ALL
from ctypes import POINTER, cast

# ===== WINDOWS =====
import win32gui
import win32con

# ===== POWERPOINT =====
from pptx import Presentation
from pptx.util import Inches, Pt
import random
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ===== GROQ =====
from dotenv import dotenv_values
from groq import Groq


class SystemAutomation:
    def __init__(self):
        self.content_cache = {}
        self.system = platform.system()
        self.data_folder = Path("Data")
        self.data_folder.mkdir(exist_ok=True)
        self.recording = False

        env = dotenv_values(".env")
        self.groq = Groq(api_key=env.get("GroqAPIKey")) if env.get("GroqAPIKey") else None

    # ==================================================
    # 📊 VOLUME (FAST + SAFE)
    # ==================================================
    def _volume_iface(self):
        try:
            devices = AudioUtilities.GetSpeakers()
            interface = devices.Activate(
                IAudioEndpointVolume._iid_, CLSCTX_ALL, None
            )
            return cast(interface, POINTER(IAudioEndpointVolume))
        except:
            return None

    def set_volume(self, level):
        try:
            level = int(level)
        except:
            return "Invalid volume level"

        level = max(0, min(100, level))
        iface = self._volume_iface()

        if iface:
            iface.SetMasterVolumeLevelScalar(level / 100.0, None)
            return f"Volume set to {level}%"

        # FINAL fallback (absolute reset)
        for _ in range(50):
            pyautogui.press("volumedown", _pause=False)
        presses = int(level / 2)   # ~2% per key
        for _ in range(presses):
            pyautogui.press("volumeup", _pause=False)

        return f"Volume set to {level}%"

    def mute_volume(self):
        iface = self._volume_iface()
        if iface:
            iface.SetMute(1, None)
        else:
            pyautogui.press("volumemute")
        return "Volume muted"

    def unmute_volume(self):
        iface = self._volume_iface()
        if iface:
            iface.SetMute(0, None)
        else:
            pyautogui.press("volumemute")
        return "Volume unmuted"

    # ==================================================
    # 🌐 SEARCH & MEDIA (NEW CAPABILITIES)
    # ==================================================
    def google_search(self, query):
        """Searches the query on Google."""
        webbrowser.open(f"https://www.google.com/search?q={query}")
        return f"Searching Google for: {query}"

    def youtube_search(self, query):
        """Searches the query on YouTube."""
        webbrowser.open(f"https://www.youtube.com/results?search_query={query}")
        return f"Searching YouTube for: {query}"

    def play_on_youtube(self, topic):
        """Uses the 'results' page with an auto-play trigger for the first video."""
        try:
            import pywhatkit
            pywhatkit.playonyt(topic)
            return f"Playing {topic} on YouTube."
        except:
            url = f"https://www.youtube.com/results?search_query={topic}"
            webbrowser.open(url)
            return f"Opening YouTube for: {topic}"

    # ==================================================
    # 🔍 APP EXISTENCE CHECK
    # ==================================================
    def _check_app_exists(self, name):
        """
        Check if app exists in common locations or is running
        Returns: (exists: bool, path: str or None)
        """
        name_lower = name.lower().strip()
        
        # Check if already running
        for proc in psutil.process_iter(['name', 'exe']):
            try:
                proc_name = proc.info['name'].lower() if proc.info['name'] else ''
                proc_exe = proc.info['exe'].lower() if proc.info['exe'] else ''
                
                if name_lower in proc_name or name_lower in proc_exe:
                    return True, proc.info['exe']
            except (psutil.NoSuchProcess, psutil.AccessDenied):
                continue
        
        # Check common install locations
        common_paths = [
            f"C:\\Program Files\\{name}\\{name}.exe",
            f"C:\\Program Files (x86)\\{name}\\{name}.exe",
            f"C:\\Program Files\\{name.capitalize()}\\{name}.exe",
            f"C:\\Program Files (x86)\\{name.capitalize()}\\{name}.exe",
        ]
        
        for path in common_paths:
            if os.path.exists(path):
                return True, path
        
        return False, None

    def _get_website_for_app(self, app_name):
        """
        Get common website URLs for popular apps
        """
        website_map = {
            'youtube': 'https://www.youtube.com',
            'gmail': 'https://mail.google.com',
            'google': 'https://www.google.com',
            'facebook': 'https://www.facebook.com',
            'twitter': 'https://twitter.com',
            'instagram': 'https://www.instagram.com',
            'linkedin': 'https://www.linkedin.com',
            'spotify': 'https://open.spotify.com',
            'netflix': 'https://www.netflix.com',
            'amazon': 'https://www.amazon.com',
            'whatsapp': 'https://web.whatsapp.com',
            'discord': 'https://discord.com/app',
            'reddit': 'https://www.reddit.com',
            'github': 'https://github.com',
            'stackoverflow': 'https://stackoverflow.com',
            'drive': 'https://drive.google.com',
            'docs': 'https://docs.google.com',
            'sheets': 'https://sheets.google.com',
            'slides': 'https://slides.google.com',
            'maps': 'https://maps.google.com',
            'translate': 'https://translate.google.com',
        }
        
        return website_map.get(app_name.lower())

    # ==================================================
    # 📦 APPS (ENHANCED WITH SMART FALLBACK)
    # ==================================================
    def _find_installed_app(self, name):
        """Optimized to only look at the main shortcut files."""
        # We only check the main 'Programs' folder, no subfolders unless necessary
        start_dirs = [
            Path(os.environ.get("PROGRAMDATA", "")) / "Microsoft/Windows/Start Menu/Programs",
            Path(os.environ.get("APPDATA", "")) / "Microsoft/Windows/Start Menu/Programs",
        ]
        
        for d in start_dirs:
            if not d.exists(): continue
            # Look only for .lnk files in the top level (Much faster than rglob)
            for file in d.glob("*.lnk"):
                if name in file.stem.lower():
                    return str(file)
        return None

    def _find_website(self, name):
        for prefix in ["https://", "https://www."]:
            for domain in ["com", "org", "net", "io"]:
                url = f"{prefix}{name}.{domain}"
                try:
                    r = requests.head(url, timeout=3, allow_redirects=True)
                    if r.status_code < 400:
                        return r.url
                except:
                    pass
        return None


    def open_app(self, name):
        name = name.lower().strip()
        
        # 1. UWP / SYSTEM APP MAP (Modern Windows Apps)
        # These don't have paths; they use URI schemes
        uwp_map = {
            "settings": "ms-settings:root",
            "calculator": "calculator:",
            "paint": "mspaint",
            "photos": "ms-photos:",
            "camera": "microsoft.windows.camera:",
            "calendar": "outlookcal:",
            "store": "ms-windows-store:",
            "microsoft store": "ms-windows-store:",
            "weather": "bingweather:",
            "notepad": "notepad",
        }

        # 2. TRADITIONAL EXE MAP (Classic Desktop Apps)
        app_map = {
            "chrome": r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
            "brave": r"C:\Program Files\BraveSoftware\Brave-Browser\Application\brave.exe",
            "microsoft word": "winword.exe",
            "word": "winword.exe",
            "microsoft powerpoint": "powerpnt.exe",
            "powerpoint": "powerpnt.exe",
            "excel": "excel.exe",
            "microsoft excel": "excel.exe",
            "vlc": r"C:\Program Files\VideoLAN\VLC\vlc.exe"
        }

        try:
            # Check UWP Map first (Settings, Paint 3D, etc.)
            if name in uwp_map:
                os.system(f"start {uwp_map[name]}")
                return f"Opened {name} system app."

            # Check Desktop App Map
            if name in app_map:
                os.startfile(app_map[name])
                return f"Opened {name} instantly."

            # 3. SHALLOW SEARCH (For everything else)
            path = self._find_installed_app(name)
            if path:
                os.startfile(path)
                return f"Opened {name}"

        except Exception as e:
            print(f"DEBUG: Error opening {name}: {e}")

        # 4. WEB FALLBACK (Only if not a system app)
        site = self._find_website(name)
        if site:
            webbrowser.open(site)
            return {
                "display":f"Opened {site}",
                "speech":f"Opened {name}"
            }

        return self.google_search(name)
    
    

    def close_app(self, name):
        """
        Enhanced close_app with better process matching
        """
        name = name.lower().strip()
        closed_count = 0
        
        try:
            # Common process name variations
            possible_names = [
                name,
                f"{name}.exe",
                name.replace(" ", ""),
                name.capitalize(),
                f"{name.capitalize()}.exe"
            ]
            
            for proc in psutil.process_iter(['name', 'exe']):
                try:
                    proc_name = proc.info['name'].lower() if proc.info['name'] else ''
                    proc_exe = proc.info['exe'].lower() if proc.info['exe'] else ''
                    
                    # Check multiple matching criteria
                    should_close = False
                    
                    # Exact name match
                    if any(pn.lower() == proc_name for pn in possible_names):
                        should_close = True
                    
                    # Name contains our search term
                    elif name in proc_name:
                        should_close = True
                    
                    # Exe path contains our search term
                    elif proc_exe and name in proc_exe:
                        should_close = True
                    
                    if should_close:
                        proc.terminate()
                        closed_count += 1
                        
                except (psutil.NoSuchProcess, psutil.AccessDenied, psutil.TimeoutExpired):
                    continue
            
            # Wait a bit for processes to close
            if closed_count > 0:
                import time
                time.sleep(0.5)
                return f"Closed {closed_count} instance(s) of {name}."
            else:
                return f"{name} is not running."
                
        except Exception as e:
            return f"Error closing {name}: {str(e)}"

    # ==================================================
    # 📝 CONTENT WRITING
    # ==================================================
    def write_content(self, topic, content_type="letter"):
        if not self.groq:
            return "Groq API not available for content generation."

        prompt = f"Write a professional {content_type} on the topic: {topic}. Keep it concise and well-structured."

        try:
            completion = self.groq.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=1024,
            )
            content = completion.choices[0].message.content.strip()

            filename = f"{content_type.capitalize()}_{topic.replace(' ', '_')}_{datetime.now():%Y%m%d_%H%M%S}.txt"
            path = self.data_folder / filename
            with open(path, 'w', encoding='utf-8') as f:
                f.write(content)

            subprocess.Popen(['notepad.exe', str(path)])

            return f"{content_type.capitalize()} on '{topic}' created and opened in Notepad."

        except Exception as e:
            return f"Error generating content: {str(e)}"

    # ==================================================
    # 📊 PRESENTATION CREATION
    # ==================================================
    
    def _get_random_theme(self):
        themes = [
            {"bg": (20, 20, 20), "title": (0, 210, 255), "text": (255, 255, 255)}, # Cyber Dark
            {"bg": (240, 240, 240), "title": (44, 62, 80), "text": (52, 73, 94)},  # Clean Professional
            {"bg": (26, 54, 104), "title": (255, 215, 0), "text": (255, 255, 255)}, # Royal Blue/Gold
            {"bg": (46, 125, 50), "title": (255, 255, 255), "text": (232, 245, 233)}, # Nature Green
            {"bg": (123, 31, 162), "title": (255, 235, 59), "text": (255, 255, 255)}  # Deep Purple
        ]
        return random.choice(themes)

    def _generate_dynamic_theme(self):
        """Generates a random high-contrast color theme."""
        # Random Dark Background
        bg_r, bg_g, bg_b = random.randint(10, 50), random.randint(10, 50), random.randint(10, 50)
        # Bright Accent for Title (High Saturation)
        acc_r, acc_g, acc_b = random.randint(150, 255), random.randint(150, 255), random.randint(150, 255)
        
        return {
            "bg": (bg_r, bg_g, bg_b),
            "title": (acc_r, acc_g, acc_b),
            "text": (240, 240, 240) # Off-white for readability
        }

    def create_presentation(self, topic):
        # AI now returns slide data AND a list of suggested fonts
        slides_data, suggested_fonts = self._generate_slides(topic)
        theme = self._generate_dynamic_theme()
        
        # Pick a random font from the AI's suggestions, or fallback if list is empty
        chosen_font = random.choice(suggested_fonts) if suggested_fonts else "Arial"
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for i, s in enumerate(slides_data):
            layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            
            # --- APPLY BACKGROUND ---
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*theme["bg"])

            # --- TITLE ---
            title_shape = slide.shapes.title
            title_shape.text = s["title"].upper()
            title_p = title_shape.text_frame.paragraphs[0]
            title_p.font.color.rgb = RGBColor(*theme["title"])
            title_p.font.bold = True
            title_p.font.size = Pt(44)
            title_p.font.name = chosen_font # Use AI-suggested font

            # --- CONTENT ---
            if i > 0 and "points" in s:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.word_wrap = True
                tf.clear()

                for point in s["points"]:
                    if not point.strip(): continue
                    p = tf.add_paragraph()
                    p.text = f"➤ {point.strip()}"
                    p.font.color.rgb = RGBColor(*theme["text"])
                    p.font.size = Pt(18)
                    p.font.name = "Segoe UI" # Standard clean font for body
                    p.space_after = Pt(12)
            
            elif i == 0:
                subtitle = slide.placeholders[1]
                subtitle.text = f"Comprehensive Analysis: {topic}\nGenerated by SARA Intelligence"
                subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme["text"])

        clean_topic = "".join([c for c in topic if c.isalnum() or c in (' ', '-')]).strip()
        path = self.data_folder / f"{clean_topic}.pptx"
        prs.save(path)
        os.startfile(path)
        return f"Created presentation on {topic} using '{chosen_font}' font."

    def _generate_slides(self, topic):
        if not self.groq: return self._fallback_slides(topic), ["Arial"]

        prompt = f"""
        Create a detailed 8-slide presentation about {topic}.
        
        First, list 5 professional Windows font names that suit this topic (e.g., Impact, Verdana, Calibri, Georgia, Segoe UI).
        Then, provide the slides.
        
        Format:
        FONTS: Font1, Font2, Font3, Font4, Font5
        
        SLIDE:
        TITLE: <Slide Title>
        POINT: <Detailed explanation sentence>
        POINT: <Detailed explanation sentence>
        
        ---
        """
        
        try:
            r = self.groq.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=3000
            )
            content = r.choices[0].message.content
            return self.parse_slides_v3(content)
        except:
            return self._fallback_slides(topic), ["Arial"]

    def parse_slides_v3(self, text):
        """Captures fonts and prevents the '...' dots issue by removing manual chunking."""
        fonts = []
        slides = []
        
        # Extract Fonts
        for line in text.splitlines():
            if line.startswith("FONTS:"):
                fonts = [f.strip() for f in line.replace("FONTS:", "").split(",")]
                break
        
        # Extract Slides
        raw_slides = text.split("---")
        for section in raw_slides:
            if "TITLE:" in section:
                lines = section.strip().splitlines()
                slide_obj = {"title": "Information", "points": []}
                for line in lines:
                    if line.startswith("TITLE:"):
                        slide_obj["title"] = line.replace("TITLE:", "").strip()
                    elif line.startswith("POINT:"):
                        # Captured fully without manual 150-char chunking
                        slide_obj["points"].append(line.replace("POINT:", "").strip())
                if slide_obj["points"]:
                    slides.append(slide_obj)
                    
        return slides, fonts
    
    def _fallback_slides(self, topic):
        return [
            {"title": "Introduction", "points": [f"What is {topic}", "Basic overview"]},
            {"title": "Background", "points": ["History", "Evolution"]},
            {"title": "Working", "points": ["How it works", "Key components"]},
            {"title": "Applications", "points": ["Real-world use", "Industries"]},
            {"title": "Advantages", "points": ["Benefits", "Efficiency"]},
            {"title": "Conclusion", "points": ["Summary", "Future scope"]},
        ]

    # ==================================================
    # 📸 SCREENSHOT & RECORDING
    # ==================================================
    def take_screenshot(self):
        filename = f"Screenshot_{datetime.now():%Y%m%d_%H%M%S}.png"
        path = self.data_folder / filename
        pyautogui.screenshot(str(path))
        os.startfile(path)
        return f"Screenshot saved."

    def start_screen_recording(self):
        options = ["System Audio", "Mic Audio", "System + Mic"]
        choice = pyautogui.confirm(text="Select Audio Source:", title="SARA Recorder", buttons=options)
        if not choice: return "Recording cancelled."

        self.recording = True
        self.temp_audio = str(self.data_folder / "temp_audio.wav")
        self.temp_video = str(self.data_folder / "temp_video.mp4")
        self.final_output = str(self.data_folder / f"Recording_{datetime.now():%Y%m%d_%H%M%S}.mp4")

        threading.Thread(target=self._video_recorder).start()
        threading.Thread(target=self._audio_recorder, args=(choice,)).start()

        return f"🔴 Recording ({choice}). Press F8 to STOP."

    def _audio_recorder(self, choice):
        import pyaudiowpatch as pyaudio
        p = pyaudio.PyAudio()
        
        if "System" in choice:
            wasapi_info = p.get_host_api_info_by_type(pyaudio.paWASAPI)
            device_index = p.get_default_output_device_info()['index']
            device = p.get_device_info_by_index(device_index)
            if not device["isLoopbackDevice"]:
                for loopback in p.get_loopback_device_info_generator():
                    if device["name"] in loopback["name"]:
                        device_index = loopback["index"]
                        break
        else:
            device_index = p.get_default_input_device_info()['index']

        device_info = p.get_device_info_by_index(device_index)
        native_rate = int(device_info.get("defaultSampleRate", 44100))

        stream = p.open(
            format=pyaudio.paInt16, 
            channels=2, 
            rate=native_rate,
            input=True, 
            input_device_index=device_index,
            frames_per_buffer=1024
        )
        
        frames = []
        while self.recording:
            frames.append(stream.read(1024))
        
        stream.stop_stream()
        stream.close()
        
        wf = wave.open(self.temp_audio, 'wb')
        wf.setnchannels(2)
        wf.setsampwidth(p.get_sample_size(pyaudio.paInt16))
        wf.setframerate(44100)
        wf.writeframes(b''.join(frames))
        wf.close()
        p.terminate()

    def _video_recorder(self):
        screen_size = tuple(pyautogui.size())
        fourcc = cv2.VideoWriter_fourcc(*"mp4v")
        out = cv2.VideoWriter(self.temp_video, fourcc, 20.0, screen_size)

        while self.recording:
            if keyboard.is_pressed("f8"):
                self.recording = False
                break
            
            img = pyautogui.screenshot()
            frame = np.array(img)
            frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            out.write(frame)

        out.release()
        self._merge_audio_video()

    def _merge_audio_video(self):
        cmd = f'ffmpeg -i "{self.temp_video}" -i "{self.temp_audio}" -c:v copy -c:a aac -strict experimental "{self.final_output}" -y'
        subprocess.run(cmd, shell=True)
        
        if os.path.exists(self.temp_video): os.remove(self.temp_video)
        if os.path.exists(self.temp_audio): os.remove(self.temp_audio)
        
        os.startfile(self.data_folder)
        print(f"✅ Success! Saved to {self.final_output}")